import chainlit as cl
from langchain_ollama import ChatOllama, OllamaEmbeddings
from pypdf import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
import hashlib
from docx import Document as DocxReader #DocxReader for no duplicated to Document
import csv
from pptx import Presentation
import pandas as pd

#launching models generation && embeddings
llm = ChatOllama(
    model = 'qwen2.5-coder:7b'
)

embeddings = OllamaEmbeddings(
    model = 'nomic-embed-text'
)

# check duplicate func
def get_file_hash(filepath):
    hasher = hashlib.md5()

    with open(filepath, "rb") as f:
        buf = f.read()
        hasher.update(buf)  

    return hasher.hexdigest()

# read_pdf_file
def parse_pdf(filename, source_name):
    reader = PdfReader(filename)
    docs = []

    
    for i, page in enumerate(reader.pages):
        page_text = page.extract_text() or ""

        if page_text.strip():
            docs.append(
                Document(
                    page_content=f"\n PAGE {i + 1} \n{page_text}",
                    metadata={
                        "source": source_name,
                        "page": i + 1,
                        "file_type": "pdf"
                    }
                )
            )

    return docs

# read_docs_file
def parse_docx(filename, source_name):
    reader = DocxReader(filename)
    docs = []

    paragraphs = []

    for para in reader.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)

    full_text = "\n".join(paragraphs)

    if full_text.strip():
        docs.append(
            Document(
                page_content= full_text,
                metadata={
                    "source": source_name,
                    "page": None,
                    'file_type': "docx"
                }
            )
        )

    return docs

# read_csv_file
def parse_csv(filename, source_name):
    docs = []
    rows = []

    try:
        with open( filename, newline = "", encoding= 'utf-8') as csvfile:
            reader = csv.reader(csvfile, delimiter= ',', quotechar= '"')
            for row in reader:
                rows.append(', '.join(row))
    except UnicodeDecodeError:
        with open( filename, newline = "", encoding= 'latin-1') as csvfile:
            reader = csv.reader(csvfile, delimiter= ',', quotechar= '"')
            for row in reader:
                rows.append(', '.join(row))
    
    full_text = "\n".join(rows)

    if full_text.strip():
        docs.append(
            Document(
                page_content=full_text,
                metadata={
                    "source": source_name,
                    "page": None,
                    "file_type": "csv"
                }
            )
        )

    return docs

# read_text_file
def parse_text(filename, source_name):
    docs = []
    
    try:
        with open(filename, 'r', encoding= 'utf-8') as file:
            full_text = file.read()
    except UnicodeDecodeError:
        with open(filename, 'r', encoding= 'latin-1') as file:
            full_text = file.read()

    if full_text.strip():
        docs.append(
            Document(
                page_content = full_text,
                metadata = {
                    "source": source_name,
                    "page": None,
                    "file_type": "txt"
                }
            )
        )

    return docs

# read_pptx_file
def parse_pptx(filename, source_name):
    presentation = Presentation(filename)
    docs = []
    
    total_slides = len(presentation.slides)

    for slide_number, slide in enumerate(presentation.slides):
        slide_text = ""

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text.strip() + "\n"

        if slide_text.strip():
            docs.append(
                Document(
                    page_content=f"Slide {slide_number + 1}:\n{slide_text}",
                    metadata = {
                        "source": source_name,
                        "page": slide_number + 1,
                        "total_slides": total_slides,
                        "file_type": "pptx"
                    }
                )
            )

    return docs

# read_excel_file
def parse_xlsx_xls(filename, source_name):
    sheets = pd.read_excel(filename, sheet_name= None)
    docs = []

    for sheet_name, df in sheets.items():
        full_text = df.to_csv(index= False)

        if full_text.strip():
            docs.append(
                Document(
                    page_content=f"Sheet: {sheet_name}\n{full_text}",
                    metadata={
                        "source": source_name,
                        "page": None,
                        "sheet": sheet_name,
                        "file_type": "excel"
                    }
                )
            )

    return docs

# check types of uploaded files
def check_files_type(element):
    if element.mime == 'application/pdf':
        return parse_pdf(element.path, element.name)
    
    elif element.mime == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        return parse_docx(element.path, element.name)
    
    elif element.mime in ["text/plain", "text/markdown"] or element.name.lower().endswith(".md"):
        return parse_text(element.path, element.name)
    
    elif element.mime == 'text/csv':
        return parse_csv(element.path, element.name)
    
    elif element.mime == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        return parse_pptx(element.path, element.name)
    
    elif element.mime in [
                            "application/vnd.ms-excel",
                            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        ]:
        return parse_xlsx_xls(element.path, element.name)
    
    else:
        return []

#split documents using chunks
def chunk_docs(docs):
    splitter = RecursiveCharacterTextSplitter(
        separators=["\n\n", "\n", ". ", " ", ""],
        chunk_size=1000,
        chunk_overlap=300
    )

    return splitter.split_documents(docs)

# build context for the model to generate answers
def build_context(results):
    context = "\n\n".join([
        f"Source: {doc.metadata.get('source', 'unknown')}\n"
        f"File type: {doc.metadata.get('file_type', 'unknown')}\n"
        f"Page/Slide: {doc.metadata.get('page', 'N/A')}\n"
        f"Total slides: {doc.metadata.get('total_slides', 'N/A')}\n"
        f"Sheet: {doc.metadata.get('sheet', 'N/A')}\n"
        f"Content:\n{doc.page_content}"
        for doc in results
    ])

    return context

# build prompt
def build_prompt(question, context, uploaded_files_text):
    prompt = f"""
        You are a helpful document question-answering assistant.

        You answer questions ONLY using the uploaded document context below.
        The user may use informal words, typos, abbreviations, or file extensions to refer to uploaded files.

        Uploaded files:
        {uploaded_files_text}

        Context:
        {context}

        User question:
        {question}

        File reference rules:
        - If the user says "pdf", answer using PDF files.
        - If the user says "ppt", "pptx", "presentation", or "powerpoint", answer using PPTX files.
        - If the user says "xls", "xlsx", "excel", "spreadsheet", or typo like "cls", they probably mean the uploaded Excel file.
        - If the user says "csv", answer using CSV files.
        - If the user says "doc", "docs", or "docx", answer using DOCX files.
        - If the user says "txt" or "md", answer using text/markdown files.
        - If the user says "uploaded file", "this file", or "the file", infer the most relevant uploaded file from the question and context.
        - If multiple files match, prefer the file type or file name mentioned in the user question.
        - Do not say a file is missing if it appears in the Uploaded files list or Context.

        Answering rules:
        - Answer directly and clearly.
        - Use the context above.
        - Mention the source file.
        - Mention page, slide, or sheet when available.
        - If the user asks how many slides/sections are in a PPTX and the context contains Total slides, use Total slides.
        - If the answer cannot be determined from the retrieved context, say what information is missing and suggest asking about a specific uploaded file name.
        """
    
    return prompt

# call prompt
def called_prompt(question, vectorstore, k, uploaded_files=None):

    if k < 1:
        return "No documents available"

    if uploaded_files is None:
        uploaded_files = []

    uploaded_files_text = "\n".join(uploaded_files) or "No uploaded file metadata available."

    k_min = min(12, k)
    results = vectorstore.similarity_search(question, k=k_min)

    context = build_context(results)

    return build_prompt(question, context, uploaded_files_text)

# model-generation
async def stream_msg(prompt):
    msg = cl.Message(content="")
    await msg.send()

    async for chunk in llm.astream(prompt):
        token = getattr(chunk, "content", "")
        if token:
            await msg.stream_token(token)

    await msg.update()

# main
@cl.on_message
async def main(message: cl.Message):

    # Case 1: User uploads document(s)
    if message.elements:
        for element in message.elements:

            uploaded_hashes = cl.user_session.get("uploaded_hashes") or set()

            file_hash = get_file_hash(element.path)

            #check Duplicate.
            if file_hash in uploaded_hashes:
                await cl.Message(
                    content=f"{element.name} was already uploaded, so I skipped it."
                ).send()
                continue

            try:
                raw_docs = check_files_type(element)
            except Exception as e:
                await cl.Message(
                    content=f"Could not read {element.name}: {str(e)}"
                ).send()
                continue

            # inappropriate files
            if not raw_docs:
                await cl.Message(
                    content = f"File {element.name} is not supported now, \
                    so please convert your necessary file into PDF, DOCX, TXT, MD, CSV, PPTX, XLS, or XLSX"
                ).send()
                continue
            
            # split documents by using chunks
            docs = chunk_docs(raw_docs)
            
            if docs:
                #get vectorstore && all_chunks
                vectorstore = cl.user_session.get("vectorstore")
                all_docs = cl.user_session.get("chunks") or []

                #check whether there is already uploaded file
                if vectorstore is None:
                    #create FAISS vectorestore
                    vectorstore = FAISS.from_documents(
                        docs,
                        embedding= embeddings
                    )

                else:
                    #add file into the existing FAISS vectorstore
                    vectorstore.add_documents(docs)

                #extend docs into all_docs
                all_docs.extend(docs)

                # save uploaded file metadata and check duplicate files to remove them
                uploaded_files = cl.user_session.get("uploaded_files") or []

                # create file_info
                file_info = (
                    f"name={element.name}, "
                    f"mime={element.mime}, "
                    f"chunks={len(docs)}"
                )

                # get data from file_info
                uploaded_files.append(file_info)

                # set_chunks & vectorstore for reuse purposes
                uploaded_hashes.add(file_hash)
                cl.user_session.set("uploaded_hashes", uploaded_hashes)
                cl.user_session.set("chunks", all_docs)
                cl.user_session.set("vectorstore", vectorstore)
                cl.user_session.set("uploaded_files", uploaded_files)

        # in cases where users upload document(s) and ask questions simultaneously
        if message.content.strip():

            question = message.content
            vectorstore = cl.user_session.get("vectorstore")
            all_docs = cl.user_session.get("chunks") or []
            uploaded_files = cl.user_session.get("uploaded_files") or []

            # call prompt func
            prompt = called_prompt(
                question,
                vectorstore,
                len(all_docs),
                uploaded_files=uploaded_files
            )

            # response msg
            await stream_msg(prompt)

        else:
            await cl.Message(
                content = "Document uploaded successfully. Now ask me a question about it."
            ).send()

        return
    else: 
        # Case 2: User asks a question
        vectorstore = cl.user_session.get('vectorstore')
        chunks = cl.user_session.get('chunks')

        # check whether there is already uploaded files before LLMs model respond to users' questions.
        if vectorstore is None:
            await cl.Message(
                content = '''Please upload a document first. Supported files: PDF, DOCX, TXT, CSV, PPTX, XLS, XLSX, MD.'''
            ).send()
            return

        # get uploaded_files data
        uploaded_files = cl.user_session.get("uploaded_files") or []
        
        # call prompt
        prompt = called_prompt(message.content, vectorstore, len(chunks), uploaded_files=uploaded_files)
        
        # response msg
        await stream_msg(prompt)